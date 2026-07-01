"""Tests for the document-generation tools (generate_proposal/ringisho/pptx/docx).

The deterministic path (proposal/ringisho) is exercised fully — no GPU/LLM, since
conftest leaves SENPAI_USE_LLM off. The general tools (pptx/docx) require a model, so
here we only assert they degrade gracefully (a clear message, no file) when it's off.
All output is redirected to a tmp dir so the committed seed is never touched.
"""
from __future__ import annotations

import pytest
from docx import Document
from pptx import Presentation

from senpai import config
from senpai.data import store
from senpai.documents import proposal, registry, ringisho
from senpai.documents.context import build_document_context
from senpai.documents.render import render_docx, render_pptx
from senpai.tools.impl import dispatch

DEAL = "D001"  # seeded dead-but-optimistic deal with real pain points + financials


@pytest.fixture(autouse=True)
def _tmp_generated(tmp_path, monkeypatch):
    """Redirect generated files to a tmp dir for every test."""
    monkeypatch.setattr(config, "GENERATED_DIR", tmp_path / "generated")
    return tmp_path / "generated"


def _pptx_text(path) -> str:
    out = []
    for slide in Presentation(str(path)).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
    return "\n".join(out)


def _docx_text(path) -> str:
    return "\n".join(p.text for p in Document(str(path)).paragraphs)


# --- render.py (pure, LLM-free) ------------------------------------------------
def test_render_pptx_slide_count(tmp_path):
    spec = {"slides": [
        {"layout": "title", "title": "T", "subtitle": "sub"},
        {"layout": "content", "title": "A", "bullets": ["x", "y"]},
        {"layout": "content", "title": "B", "bullets": ["z"]},
    ]}
    p = render_pptx(spec, tmp_path / "d.pptx")
    prs = Presentation(str(p))
    assert len(prs.slides) == 3


def test_render_docx_headings(tmp_path):
    spec = {"title": "Doc", "sections": [
        {"heading": "One", "body": ["para a", "- bullet"]},
        {"heading": "Two", "body": ["para b"]},
    ]}
    p = render_docx(spec, tmp_path / "d.docx")
    text = _docx_text(p)
    assert "One" in text and "Two" in text and "bullet" in text


# --- context grounding ---------------------------------------------------------
def test_context_numbers_match_store():
    ctx = build_document_context(DEAL)
    d = store.get_deal(DEAL)
    assert ctx is not None
    assert ctx.financials["investment"] == d["total_order_amount"]
    assert ctx.customer == store.customer_name(d["customer_id"])
    assert ctx.pain_points  # real customer_challenge values exist for this deal


def test_context_unknown_deal_is_none():
    assert build_document_context("ZZZ") is None


# --- proposal (PPTX) -----------------------------------------------------------
def test_proposal_arc_and_grounded():
    res = proposal.generate(DEAL)
    assert res is not None
    path, ctx, spec = res
    prs = Presentation(str(path))
    # Full proposal arc: 表紙 → 背景 → 課題 → ソリューション → 投資対効果 → 次のステップ.
    assert len(prs.slides) == 6
    assert len(spec["slides"]) == 6
    text = _pptx_text(path)
    assert ctx.customer in text                       # title slide names the customer
    assert f"{ctx.financials['investment']:,}" in text  # ROI slide carries the real ¥ (D001 has no quote)


# --- ringisho (DOCX) -----------------------------------------------------------
def test_ringisho_headings_and_amount():
    res = ringisho.generate(DEAL)
    assert res is not None
    path, ctx = res
    text = _docx_text(path)
    for heading in ("稟議書", "背景・課題", "提案内容", "投資額と効果", "結論・承認依頼"):
        assert heading in text
    assert f"{ctx.financials['investment']:,}" in text


# --- pptx generates directly (no confirm gate) ---------------------------------
def test_proposal_tool_generates_directly(_tmp_generated):
    gen = _tmp_generated
    # PPTX proposals build in one round — no preview/confirm step.
    out = dispatch("generate_proposal", {"deal_id": DEAL})
    assert "生成しました" in out
    assert "プレビュー" not in out
    assert len(list(gen.glob("*.pptx"))) == 1                 # file written on the first call


def test_ringisho_tool_writes_docx(_tmp_generated):
    dispatch("generate_ringisho", {"deal_id": DEAL, "confirm": True})
    assert len(list(_tmp_generated.glob("*.docx"))) == 1


# --- general tools need the model ----------------------------------------------
def test_general_tools_need_model(_tmp_generated):
    msg = dispatch("generate_pptx", {"prompt": "GTA 6", "confirm": True})
    assert "モデル" in msg                                     # "needs the model"
    msg2 = dispatch("generate_docx", {"prompt": "security training"})
    assert "モデル" in msg2
    assert not _tmp_generated.exists() or not list(_tmp_generated.iterdir())  # no file


# --- grounding: general tools ground on conversation + workspace, not just CRM --
def test_gather_grounding_uses_conversation_and_workspace(tmp_path, monkeypatch):
    """A 'proposal for <company>' where the company lives in the rep's local files
    (and was discussed earlier) must ground on that file/conversation — and must NOT
    inject an unrelated fuzzy CRM customer (the wrong-company-name hallucination).
    Uses a hermetic workspace so it never depends on the configured WORKSPACE_ROOT."""
    from senpai import config
    from senpai.tools import conversation as conv
    from senpai.tools import impl

    ws = tmp_path / "workspace"
    ws.mkdir()
    (ws / "murata_printing_display_quote.txt").write_text(
        "有限会社村田印刷 様\n27インチモニター × 4台: ¥204,000\n", encoding="utf-8")
    monkeypatch.setattr(config, "WORKSPACE_ROOT", ws)

    conv.set_conversation([
        {"role": "system", "content": "sys"},
        {"role": "user", "content": "村田印刷にいくら見積もった？"},
        {"role": "tool", "content": "ワークスペース文書: 有限会社村田印刷 ¥204,000"},
        {"role": "assistant", "content": "村田印刷への見積もりは¥204,000です。"},
        {"role": "user", "content": "make a proposal ppt for Murata Printing"},
    ])
    try:
        g = impl._gather_grounding("make a proposal ppt for Murata Printing 村田印刷",
                                   customer="", use_web=False)
    finally:
        conv.set_conversation(None)
    assert "村田印刷" in g                       # the referenced entity is grounded
    assert "204,000" in g                        # its real figure, from file/conversation
    assert "松田" not in g                        # no unrelated fuzzy CRM customer
    assert "【社内データ】" not in g              # CRM suppressed when workspace matched


def test_gather_grounding_junk_gated_and_crm_fallback(tmp_path, monkeypatch):
    """An unrelated topic pulls no workspace junk; a real CRM customer still grounds.
    The workspace root is pointed at an empty dir so the assertion is deterministic
    regardless of what real files exist under the configured WORKSPACE_ROOT."""
    from senpai import config
    from senpai.tools import conversation as conv
    from senpai.tools import impl

    empty = tmp_path / "empty_workspace"
    empty.mkdir()
    monkeypatch.setattr(config, "WORKSPACE_ROOT", empty)  # sandbox re-reads config each call

    conv.set_conversation(None)
    assert impl._workspace_grounding("best gaming laptops under 1000000 yen") == ""
    # A named CRM customer with no local file still injects internal records.
    g = impl._gather_grounding("藤本食品の提案書", customer="", use_web=False)
    assert "【社内データ】" in g


# --- registry + isolation ------------------------------------------------------
def test_registry_records_for_download(_tmp_generated):
    dispatch("generate_proposal", {"deal_id": DEAL, "confirm": True})
    # _DOCS is a process-global registry; take the most recent proposal record.
    rec = next(r for r in reversed(list(registry._DOCS.values())) if r["kind"] == "proposal")
    assert registry.get(rec["doc_id"]) is rec
    assert str(_tmp_generated) in rec["path"]                  # under tmp, not the seed


def test_seed_dir_not_written():
    proposal.generate(DEAL)
    ringisho.generate(DEAL)
    # generated files land under the (tmp) GENERATED_DIR, never the committed seed
    assert not list(config.SEED_DIR.glob("*.pptx"))
    assert not list(config.SEED_DIR.glob("*.docx"))
