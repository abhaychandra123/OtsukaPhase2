"""Shared, LLM-free renderer — the only place python-pptx / python-docx are used.

Both the grounded tools (proposal/ringisho) and the general tools (pptx/docx) build
a normalized *spec* and hand it here. Keeping all rendering in one module means the
binary-format code is written and tested once.

Spec shapes
-----------
deck_spec (PPTX):
    {"slides": [
        {"layout": "title", "title": str, "subtitle": str},
        {"layout": "content", "title": str, "bullets": [str], "notes": str},
        ...
    ]}
doc_spec (DOCX):
    {"title": str, "subtitle": str,
     "sections": [{"heading": str, "body": [str]}, ...]}

`python-pptx` / `python-docx` are imported lazily so a missing lib can never break
the import of senpai.tools (mirrors senpai/tools/gcal.py).
"""
from __future__ import annotations

import re
from datetime import datetime
from pathlib import Path

from senpai import config


def output_path(kind: str, slug: str, ext: str) -> Path:
    """A unique, safe path under GENERATED_DIR, e.g. proposal_D001_20260616-1430.pptx.
    Creates the dir on first use."""
    config.GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    safe = re.sub(r"[^A-Za-z0-9_-]+", "_", (slug or "doc")).strip("_") or "doc"
    stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    return config.GENERATED_DIR / f"{kind}_{safe}_{stamp}.{ext}"


# Layouts we populate, mapped to the Otsuka template's layout names (and the
# blank default's, as a fallback). The template is hand-built: its "Title Only"
# layout (タイトルのみ) carries NO body placeholder, so the content layout MUST be
# resolved by name — blindly using slide_layouts[1] would silently drop bullets.
_TITLE_LAYOUT_NAMES = ("タイトル スライド", "Title Slide")
_CONTENT_LAYOUT_NAMES = ("タイトルとコンテンツ", "Title and Content")


def _layout(prs, names, fallback_idx):
    """Pick a slide layout by name (Otsuka template / Office default), else by index."""
    by_name = {layout.name: layout for layout in prs.slide_layouts}
    for name in names:
        if name in by_name:
            return by_name[name]
    return prs.slide_layouts[fallback_idx]


def render_pptx(deck_spec: dict, path: Path) -> Path:
    """Render a deck spec to a .pptx file at `path`. Returns the path.

    Opens the committed Otsuka brand template (config.PPTX_TEMPLATE_PATH) as the
    base so the deck inherits its masters/layouts/theme; falls back to python-pptx's
    blank default if the template is missing (e.g. in CI without the asset).
    """
    from pptx import Presentation
    from pptx.util import Pt

    tmpl = config.PPTX_TEMPLATE_PATH
    prs = Presentation(str(tmpl)) if tmpl.exists() else Presentation()
    title_layout = _layout(prs, _TITLE_LAYOUT_NAMES, 0)       # Title Slide
    content_layout = _layout(prs, _CONTENT_LAYOUT_NAMES, 1)   # Title and Content

    slides = deck_spec.get("slides") or []
    if not slides:  # never produce an empty deck
        slides = [{"layout": "title", "title": deck_spec.get("title", "Document"),
                   "subtitle": deck_spec.get("subtitle", "")}]

    for spec in slides:
        is_title = spec.get("layout") == "title"
        slide = prs.slides.add_slide(title_layout if is_title else content_layout)

        if slide.shapes.title is not None:
            slide.shapes.title.text = str(spec.get("title", ""))

        if is_title:
            # subtitle placeholder (idx 1) on the title layout
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    ph.text = str(spec.get("subtitle", ""))
                    break
            continue

        bullets = [str(b) for b in (spec.get("bullets") or []) if str(b).strip()]
        body = next((ph for ph in slide.placeholders
                     if ph.placeholder_format.idx == 1), None)
        if body is not None and bullets:
            tf = body.text_frame
            tf.word_wrap = True
            tf.text = bullets[0]
            for b in bullets[1:]:
                p = tf.add_paragraph()
                p.text = b
            for p in tf.paragraphs:  # readable on a single content slide
                for run in p.runs:
                    run.font.size = Pt(18)

        notes = str(spec.get("notes", "") or "").strip()
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(str(path))
    return path


def render_docx(doc_spec: dict, path: Path) -> Path:
    """Render a doc spec to a .docx file at `path`. Returns the path."""
    from docx import Document

    doc = Document()
    title = str(doc_spec.get("title", "") or "")
    if title:
        doc.add_heading(title, level=0)
    subtitle = str(doc_spec.get("subtitle", "") or "").strip()
    if subtitle:
        doc.add_paragraph(subtitle).italic = True

    for section in doc_spec.get("sections") or []:
        heading = str(section.get("heading", "") or "").strip()
        if heading:
            doc.add_heading(heading, level=1)
        for para in section.get("body") or []:
            text = str(para).strip()
            if not text:
                continue
            # A leading "- " or "・" renders as a bullet list item.
            if text[:2] in ("- ", "・") or text.startswith("• "):
                doc.add_paragraph(text.lstrip("-・• ").strip(), style="List Bullet")
            else:
                doc.add_paragraph(text)

    doc.save(str(path))
    return path
